#!/usr/bin/env python3
"""
Универсальный конвертер документов в Markdown.
Объединяет возможности markitdown + существующие инструменты (mammoth, python-pptx, PyMuPDF).

Стратегия fallback для каждого формата:
- PDF: PyMuPDF (лучше для научных статей) → markitdown (универсальный)
- DOCX: mammoth (сохраняет структуру) → markitdown (универсальный)
- PPTX: python-pptx (извлекает заметки) → markitdown (универсальный)
- Excel, Images, Audio, Video, YouTube: markitdown (новые возможности)

Использование:
    python unified_converter.py --input document.pdf --output output.md
    python unified_converter.py --input document.docx --output output.md --format docx
    python unified_converter.py --input presentation.pptx --output output.md --media-dir images/
    python unified_converter.py --input audio.mp3 --output transcript.md --format audio
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Optional

# Импорты для format handlers
try:
    import markitdown
    MARKITDOWN_AVAILABLE = True
except ImportError:
    MARKITDOWN_AVAILABLE = False
    print("Предупреждение: markitdown не установлен. Некоторые форматы будут недоступны.", file=sys.stderr)

try:
    import mammoth
    MAMMOTH_AVAILABLE = True
except ImportError:
    MAMMOTH_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False


class ConversionResult:
    """Результат конвертации"""
    def __init__(self, markdown: str, media_files: list[Path] = None, method: str = "unknown"):
        self.markdown = markdown
        self.media_files = media_files or []
        self.method = method


class UnifiedConverter:
    """Универсальный конвертер документов"""
    
    def __init__(self, input_file: Path, output_file: Path, media_dir: Optional[Path] = None):
        self.input_file = input_file.resolve()
        self.output_file = output_file.resolve()
        
        if media_dir is None:
            self.media_dir = self.output_file.parent / f"{self.output_file.stem}_media"
        else:
            self.media_dir = media_dir.resolve()
        
        self.media_dir.mkdir(parents=True, exist_ok=True)
        self.output_file.parent.mkdir(parents=True, exist_ok=True)
    
    def convert(self, format_hint: Optional[str] = None) -> ConversionResult:
        """
        Конвертировать документ в Markdown.
        
        Args:
            format_hint: Подсказка формата ('pdf', 'docx', 'pptx', 'excel', 'image', 'audio', 'video', 'youtube')
        
        Returns:
            ConversionResult с markdown текстом и списком медиа-файлов
        """
        if not self.input_file.exists():
            raise FileNotFoundError(f"Входной файл не найден: {self.input_file}")
        
        # Определить формат
        if format_hint:
            format_type = format_hint.lower()
        else:
            format_type = self._detect_format()
        
        # Выбрать обработчик
        handler_map = {
            'pdf': self._convert_pdf,
            'docx': self._convert_docx,
            'pptx': self._convert_pptx,
            'excel': self._convert_excel,
            'image': self._convert_image,
            'audio': self._convert_audio,
            'video': self._convert_video,
            'youtube': self._convert_youtube,
        }
        
        handler = handler_map.get(format_type)
        if not handler:
            # Fallback на markitdown для неизвестных форматов
            return self._convert_with_markitdown()
        
        return handler()
    
    def _detect_format(self) -> str:
        """Определить формат файла по расширению"""
        ext = self.input_file.suffix.lower()
        
        format_map = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.doc': 'docx',
            '.pptx': 'pptx',
            '.ppt': 'pptx',
            '.ppsx': 'pptx',
            '.xlsx': 'excel',
            '.xls': 'excel',
            '.csv': 'excel',
            '.png': 'image',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.gif': 'image',
            '.bmp': 'image',
            '.tiff': 'image',
            '.webp': 'image',
            '.mp3': 'audio',
            '.wav': 'audio',
            '.m4a': 'audio',
            '.flac': 'audio',
            '.mp4': 'video',
            '.avi': 'video',
            '.mov': 'video',
            '.mkv': 'video',
        }
        
        return format_map.get(ext, 'unknown')
    
    def _convert_pdf(self) -> ConversionResult:
        """
        Конвертировать PDF в Markdown.
        Стратегия: PyMuPDF (лучше для научных статей) → markitdown (fallback)
        """
        # Попробовать PyMuPDF
        if PYMUPDF_AVAILABLE:
            try:
                return self._convert_pdf_pymupdf()
            except Exception as e:
                print(f"PyMuPDF не удалось: {e}. Пробую markitdown...", file=sys.stderr)
        
        # Fallback на markitdown
        return self._convert_with_markitdown()
    
    def _convert_pdf_pymupdf(self) -> ConversionResult:
        """Конвертировать PDF используя PyMuPDF"""
        doc = fitz.open(self.input_file)
        lines = [f"<!-- Конвертировано из {self.input_file.name} с помощью PyMuPDF -->\n\n"]
        media_files = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            lines.append(f"\n## Страница {page_num + 1}\n\n")
            
            # Извлечь текст
            text = page.get_text()
            if text.strip():
                lines.append(text)
                lines.append("\n\n")
            
            # Извлечь изображения
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    img_filename = f"page{page_num + 1:03d}_img{img_index + 1:03d}.{image_ext}"
                    img_path = self.media_dir / img_filename
                    img_path.write_bytes(image_bytes)
                    media_files.append(img_path)
                    
                    rel_path = self._relative_path(img_path)
                    lines.append(f"![Изображение {img_index + 1}]({rel_path})\n\n")
                except Exception as e:
                    print(f"Предупреждение: не удалось извлечь изображение {img_index + 1} со страницы {page_num + 1}: {e}", file=sys.stderr)
        
        doc.close()
        markdown = "".join(lines)
        return ConversionResult(markdown, media_files, method="PyMuPDF")
    
    def _convert_docx(self) -> ConversionResult:
        """
        Конвертировать DOCX в Markdown.
        Стратегия: mammoth (сохраняет структуру) → markitdown (fallback)
        """
        # Попробовать mammoth
        if MAMMOTH_AVAILABLE:
            try:
                return self._convert_docx_mammoth()
            except Exception as e:
                print(f"mammoth не удалось: {e}. Пробую markitdown...", file=sys.stderr)
        
        # Fallback на markitdown
        return self._convert_with_markitdown()
    
    def _convert_docx_mammoth(self) -> ConversionResult:
        """Конвертировать DOCX используя mammoth"""
        media_files = []
        counter = [0]
        
        def save_image(image):
            counter[0] += 1
            content_type = getattr(image, "content_type", "") or ""
            ext = self._extension_for_content_type(content_type)
            filename = f"img_{counter[0]:04d}.{ext}"
            out_path = self.media_dir / filename
            
            try:
                with image.open() as f:
                    out_path.write_bytes(f.read())
                media_files.append(out_path)
            except Exception as e:
                print(f"Предупреждение: не удалось извлечь изображение ({filename}): {e}", file=sys.stderr)
                return {"src": "", "alt": ""}
            
            rel_path = self._relative_path(out_path)
            alt = getattr(image, "alt_text", None) or ""
            return {"src": rel_path, "alt": alt}
        
        image_converter = mammoth.images.img_element(save_image)
        
        with self.input_file.open("rb") as docx_file:
            result = mammoth.convert_to_markdown(docx_file, convert_image=image_converter)
        
        for msg in result.messages:
            text = getattr(msg, "message", str(msg))
            typ = getattr(msg, "type", "message")
            print(f"mammoth [{typ}]: {text}", file=sys.stderr)
        
        text = (result.value or "").strip()
        header = f"<!-- Конвертировано из {self.input_file.name} с помощью mammoth -->\n\n"
        markdown = header + text + ("\n" if text else "")
        
        return ConversionResult(markdown, media_files, method="mammoth")
    
    def _convert_pptx(self) -> ConversionResult:
        """
        Конвертировать PPTX в Markdown.
        Стратегия: python-pptx (извлекает заметки) → markitdown (fallback)
        """
        # Попробовать python-pptx
        if PPTX_AVAILABLE:
            try:
                return self._convert_pptx_pythonpptx()
            except Exception as e:
                print(f"python-pptx не удалось: {e}. Пробую markitdown...", file=sys.stderr)
        
        # Fallback на markitdown
        return self._convert_with_markitdown()
    
    def _convert_pptx_pythonpptx(self) -> ConversionResult:
        """Конвертировать PPTX используя python-pptx"""
        prs = Presentation(str(self.input_file))
        lines = [f"<!-- Конвертировано из {self.input_file.name} с помощью python-pptx -->\n\n"]
        media_files = []
        img_counter = [0]
        
        for sn, slide in enumerate(prs.slides, start=1):
            lines.append(f"\n## Слайд {sn}\n\n")
            
            for shape in self._walk_shapes(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image
                        ext = (img.ext or "png").lower()
                        if ext == "jpeg":
                            ext = "jpg"
                        img_counter[0] += 1
                        fname = f"slide{sn:02d}_img{img_counter[0]:04d}.{ext}"
                        out_img = self.media_dir / fname
                        out_img.write_bytes(img.blob)
                        media_files.append(out_img)
                        
                        rel_path = self._relative_path(out_img)
                        lines.append(f"![Изображение]({rel_path})\n\n")
                    except Exception as e:
                        print(f"Предупреждение: не удалось извлечь изображение со слайда {sn}: {e}", file=sys.stderr)
                    continue
                
                block = self._shape_text(shape)
                if block:
                    lines.append(block)
                    lines.append("\n\n")
            
            # Извлечь заметки докладчика
            try:
                nf = slide.notes_slide.notes_text_frame
                note_txt = (nf.text or "").strip() if nf is not None else ""
                if note_txt:
                    lines.append("\n**Заметки докладчика:**\n\n")
                    lines.append(note_txt)
                    lines.append("\n\n")
            except (AttributeError, ValueError):
                pass
        
        markdown = "".join(lines).rstrip() + "\n"
        return ConversionResult(markdown, media_files, method="python-pptx")
    
    def _convert_excel(self) -> ConversionResult:
        """Конвертировать Excel в Markdown используя markitdown"""
        return self._convert_with_markitdown()
    
    def _convert_image(self) -> ConversionResult:
        """Конвертировать изображение в Markdown с OCR используя markitdown"""
        return self._convert_with_markitdown()
    
    def _convert_audio(self) -> ConversionResult:
        """Конвертировать аудио в Markdown с транскрипцией используя markitdown"""
        return self._convert_with_markitdown()
    
    def _convert_video(self) -> ConversionResult:
        """Конвертировать видео в Markdown используя markitdown"""
        return self._convert_with_markitdown()
    
    def _convert_youtube(self) -> ConversionResult:
        """Конвертировать YouTube URL в Markdown используя markitdown"""
        return self._convert_with_markitdown()
    
    def _convert_with_markitdown(self) -> ConversionResult:
        """Конвертировать используя markitdown (универсальный fallback)"""
        if not MARKITDOWN_AVAILABLE:
            raise ImportError("markitdown не установлен. Установите: pip install markitdown")
        
        md = markitdown.MarkItDown()
        result = md.convert(str(self.input_file))
        
        markdown = result.text_content
        header = f"<!-- Конвертировано из {self.input_file.name} с помощью markitdown -->\n\n"
        markdown = header + markdown
        
        return ConversionResult(markdown, [], method="markitdown")
    
    def _walk_shapes(self, shapes):
        """Рекурсивно обойти все shapes в PPTX"""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._walk_shapes(shape.shapes)
            else:
                yield shape
    
    def _shape_text(self, shape) -> str:
        """Извлечь текст из shape"""
        if getattr(shape, "has_text_frame", False):
            t = (shape.text_frame.text or "").strip()
            return t
        if getattr(shape, "has_table", False):
            rows = []
            for row in shape.table.rows:
                cells = []
                for cell in row.cells:
                    cells.append((cell.text or "").strip().replace("\n", " "))
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                return "\n".join(rows)
        return ""
    
    def _extension_for_content_type(self, content_type: str) -> str:
        """Определить расширение файла по content type"""
        subtype = (content_type or "").split("/")[-1].lower().strip()
        if not subtype or subtype == "octet-stream":
            return "bin"
        if subtype == "jpeg":
            return "jpg"
        return subtype[:12]
    
    def _relative_path(self, file_path: Path) -> str:
        """Получить относительный путь от output_file к file_path"""
        try:
            return file_path.relative_to(self.output_file.parent).as_posix()
        except ValueError:
            return file_path.as_posix()
    
    def save(self, result: ConversionResult) -> None:
        """Сохранить результат конвертации"""
        self.output_file.write_text(result.markdown, encoding="utf-8")
        print(f"✓ Конвертировано: {self.input_file.name} → {self.output_file}")
        print(f"  Метод: {result.method}")
        if result.media_files:
            print(f"  Медиа-файлов: {len(result.media_files)}")
            print(f"  Директория: {self.media_dir}")


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Универсальный конвертер документов в Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Примеры использования:
  %(prog)s --input document.pdf --output output.md
  %(prog)s --input document.docx --output output.md --format docx
  %(prog)s --input presentation.pptx --output output.md --media-dir images/
  %(prog)s --input audio.mp3 --output transcript.md --format audio

Поддерживаемые форматы:
  - PDF (PyMuPDF → markitdown)
  - DOCX (mammoth → markitdown)
  - PPTX (python-pptx → markitdown)
  - Excel (markitdown)
  - Images с OCR (markitdown)
  - Audio с транскрипцией (markitdown)
  - Video (markitdown)
  - YouTube (markitdown)
        """
    )
    
    parser.add_argument(
        "-i", "--input",
        required=True,
        type=Path,
        help="Входной файл"
    )
    parser.add_argument(
        "-o", "--output",
        required=True,
        type=Path,
        help="Выходной Markdown файл"
    )
    parser.add_argument(
        "--format",
        choices=['pdf', 'docx', 'pptx', 'excel', 'image', 'audio', 'video', 'youtube'],
        help="Подсказка формата (автоопределение по расширению, если не указано)"
    )
    parser.add_argument(
        "--media-dir",
        type=Path,
        help="Директория для медиа-файлов (по умолчанию: {output_stem}_media/)"
    )
    
    args = parser.parse_args()
    
    try:
        converter = UnifiedConverter(args.input, args.output, args.media_dir)
        result = converter.convert(args.format)
        converter.save(result)
        return 0
    except Exception as e:
        print(f"Ошибка: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())

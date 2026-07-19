#!/usr/bin/env python3
"""
Пример создания русскоязычной презентации из научной статьи
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import fitz  # PyMuPDF
import json
from pathlib import Path


class Paper2PptRu:
    """Конвертер научных статей в русскоязычные презентации PPTX"""
    
    def __init__(self, pdf_path: str, output_dir: str = "paper2ppt_output"):
        self.pdf_path = pdf_path
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        
        self.figures_dir = self.output_dir / "figures"
        self.figures_dir.mkdir(exist_ok=True)
        
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)  # 16:9
        self.prs.slide_height = Inches(5.625)
        
        self.metadata = {}
        
    def extract_text_from_pdf(self) -> str:
        """Извлечение текста из PDF"""
        doc = fitz.open(self.pdf_path)
        full_text = ""
        
        for page in doc:
            full_text += page.get_text()
        
        doc.close()
        return full_text
    
    def extract_figures_from_pdf(self) -> list:
        """Извлечение рисунков из PDF"""
        doc = fitz.open(self.pdf_path)
        figure_paths = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images()
            
            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                figure_path = self.figures_dir / f"figure_{page_num}_{img_index}.{image_ext}"
                
                with open(figure_path, "wb") as f:
                    f.write(image_bytes)
                
                figure_paths.append(str(figure_path))
        
        doc.close()
        return figure_paths
    
    def add_title_slide(self, title: str, authors: str, journal: str, year: int, doi: str = ""):
        """Добавить титульный слайд"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        subtitle_text = f"{authors}\n{journal}, {year}"
        if doi:
            subtitle_text += f"\nDOI: {doi}"
        
        subtitle_shape.text = subtitle_text
        
        # Стиль
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
    
    def add_content_slide(self, title: str, content_items: list, notes: str = ""):
        """Добавить слайд с контентом"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Добавить контент
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(22)
            p.level = 0
            p.space_before = Pt(6)
        
        # Заметки докладчика
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def add_figure_slide(self, title: str, figure_path: str, caption: str = "", notes: str = ""):
        """Добавить слайд с рисунком"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Заголовок
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Рисунок
        left = Inches(1.5)
        top = Inches(1.5)
        height = Inches(3.5)
        
        try:
            pic = slide.shapes.add_picture(figure_path, left, top, height=height)
        except Exception as e:
            print(f"Ошибка добавления рисунка: {e}")
        
        # Подпись
        if caption:
            left = Inches(1.5)
            top = Inches(5.2)
            width = Inches(7)
            height = Inches(0.3)
            
            caption_box = slide.shapes.add_textbox(left, top, width, height)
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_frame.paragraphs[0].font.size = Pt(16)
            caption_frame.paragraphs[0].font.italic = True
        
        # Заметки докладчика
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def create_presentation(self, paper_data: dict):
        """Создать полную презентацию"""
        
        # 1. Титульный слайд
        self.add_title_slide(
            title=paper_data['title'],
            authors=paper_data['authors'],
            journal=paper_data['journal'],
            year=paper_data['year'],
            doi=paper_data.get('doi', '')
        )
        
        # 2. Актуальность
        self.add_content_slide(
            title="Актуальность исследования",
            content_items=[
                "Проблема X остается нерешенной",
                "Существующие методы имеют ограничения",
                "Необходим новый подход"
            ],
            notes="Актуальность нашего исследования обусловлена тем, что проблема X затрагивает значительную часть населения. Существующие методы, несмотря на их широкое применение, демонстрируют недостаточную эффективность. Поэтому разработка нового подхода является важной научной задачей."
        )
        
        # 3. Цель и задачи
        self.add_content_slide(
            title="Цель и задачи исследования",
            content_items=[
                "Цель: разработать новый метод X для решения проблемы Y",
                "Задачи:",
                "  1. Провести анализ существующих подходов",
                "  2. Разработать математическую модель",
                "  3. Экспериментально валидировать модель",
                "  4. Сравнить с современными аналогами"
            ],
            notes="Основная цель нашего исследования — разработать новый метод X для решения проблемы Y. Для достижения этой цели мы поставили четыре задачи. Первая — провести систематический анализ существующих подходов. Вторая — разработать математическую модель процесса. Третья — экспериментально валидировать модель на реальных данных. Четвертая — сравнить наш метод с современными аналогами."
        )
        
        # 4. Методы
        self.add_content_slide(
            title="Материалы и методы",
            content_items=[
                "Объект исследования: клеточная линия HeLa",
                "Методы:",
                "  • Иммунофлуоресцентная микроскопия",
                "  • Вестерн-блоттинг",
                "  • ПЦР в реальном времени",
                "Статистический анализ: t-критерий Стьюдента, p < 0,05"
            ],
            notes="В качестве объекта исследования мы использовали клеточную линию HeLa. Применяли три основных метода: иммунофлуоресцентную микроскопию для визуализации белков, вестерн-блоттинг для количественной оценки экспрессии, и ПЦР в реальном времени для анализа уровня мРНК. Статистическую значимость различий оценивали с помощью t-критерия Стьюдента при уровне значимости p < 0,05."
        )
        
        # 5. Результаты (с рисунком)
        figures = self.extract_figures_from_pdf()
        if figures:
            self.add_figure_slide(
                title="Результат 1: Экспрессия белка X увеличивается в 3 раза",
                figure_path=figures[0],
                caption="Рисунок 1. Иммунофлуоресцентное окрашивание белка X",
                notes="На этом слайде представлены результаты иммунофлуоресцентного окрашивания. Как видно на рисунке, экспрессия белка X в экспериментальной группе увеличивается в три раза по сравнению с контролем. Это подтверждает нашу гипотезу о том, что белок X играет ключевую роль в исследуемом процессе."
            )
        
        # 6. Обсуждение
        self.add_content_slide(
            title="Обсуждение результатов",
            content_items=[
                "Полученные результаты согласуются с данными литературы",
                "Новизна: впервые показана роль белка X в процессе Y",
                "Ограничения: исследование проведено in vitro",
                "Перспективы: необходимы исследования in vivo"
            ],
            notes="Обсуждая полученные результаты, следует отметить, что они хорошо согласуются с данными литературы. Новизна нашей работы заключается в том, что впервые показана роль белка X в процессе Y. Вместе с тем, наше исследование имеет ограничения — оно проведено только in vitro. В перспективе необходимо провести исследования in vivo для подтверждения полученных результатов."
        )
        
        # 7. Выводы
        self.add_content_slide(
            title="Выводы",
            content_items=[
                "1. Разработан новый метод X с точностью 95%",
                "2. Показана роль белка X в процессе Y",
                "3. Метод превосходит существующие аналоги",
                "4. Результаты открывают перспективы для клинического применения"
            ],
            notes="В заключение сформулируем основные выводы нашего исследования. Во-первых, мы разработали новый метод X, который демонстрирует точность 95 процентов. Во-вторых, впервые показана роль белка X в процессе Y. В-третьих, наш метод превосходит существующие аналоги по ключевым параметрам. В-четвертых, полученные результаты открывают перспективы для клинического применения."
        )
        
        # 8. Благодарности
        self.add_content_slide(
            title="Благодарности",
            content_items=[
                "Финансирование: РНФ, грант № 00-00-00000",
                "Соавторы: Иванов И.И., Петров П.П., Сидоров С.С.",
                "Контакты: author@university.ru"
            ],
            notes=""
        )
        
        # Сохранить метаданные
        self.metadata = {
            'title': paper_data['title'],
            'authors': paper_data['authors'],
            'journal': paper_data['journal'],
            'year': paper_data['year'],
            'doi': paper_data.get('doi', ''),
            'paper_type': paper_data.get('paper_type', 'experimental'),
            'slides_count': len(self.prs.slides),
            'figures_used': len(figures),
            'created_date': '2026-05-12',
            'language': 'ru'
        }
        
        # Сохранить презентацию
        output_path = self.output_dir / "presentation_ru.pptx"
        self.prs.save(str(output_path))
        
        # Сохранить метаданные
        metadata_path = self.output_dir / "metadata.json"
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(self.metadata, f, ensure_ascii=False, indent=2)
        
        print(f"✅ Презентация создана: {output_path}")
        print(f"✅ Слайдов: {self.metadata['slides_count']}")
        print(f"✅ Рисунков: {self.metadata['figures_used']}")
        
        return str(output_path)


def main():
    """Пример использования"""
    
    # Данные статьи
    paper_data = {
        'title': 'Роль белка X в регуляции клеточного цикла',
        'authors': 'Иванов И.И., Петров П.П., Сидоров С.С.',
        'journal': 'Молекулярная биология',
        'year': 2024,
        'doi': '10.1134/S0026893324010000',
        'paper_type': 'experimental'
    }
    
    # Создать презентацию
    converter = Paper2PptRu(
        pdf_path="paper.pdf",
        output_dir="paper2ppt_output"
    )
    
    output_path = converter.create_presentation(paper_data)
    print(f"\n🎉 Готово! Презентация сохранена: {output_path}")


if __name__ == "__main__":
    main()

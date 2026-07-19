#!/usr/bin/env python3
"""
Утилиты для верификации качества презентаций
"""

from pptx import Presentation
from pathlib import Path
import re


class PresentationVerifier:
    """Верификатор качества презентаций"""
    
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
        self.errors = []
        self.warnings = []
        self.info = []
    
    def verify_structure(self) -> bool:
        """Проверка структуры презентации"""
        print("🔍 Проверка структуры...")
        
        # Минимальное количество слайдов
        if len(self.prs.slides) < 5:
            self.errors.append(f"Недостаточно слайдов: {len(self.prs.slides)} (минимум 5)")
        else:
            self.info.append(f"✅ Количество слайдов: {len(self.prs.slides)}")
        
        # Максимальное количество слайдов
        if len(self.prs.slides) > 30:
            self.warnings.append(f"Слишком много слайдов: {len(self.prs.slides)} (рекомендуется до 25)")
        
        # Титульный слайд должен иметь заголовок
        if len(self.prs.slides) > 0:
            first_slide = self.prs.slides[0]
            has_title = False
            for shape in first_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_title = True
                    break
            
            if not has_title:
                self.errors.append("Титульный слайд не имеет заголовка")
            else:
                self.info.append("✅ Титульный слайд имеет заголовок")
        
        return len(self.errors) == 0
    
    def verify_language(self) -> bool:
        """Проверка языка (должен быть русский)"""
        print("🔍 Проверка языка...")
        
        total_chars = 0
        cyrillic_chars = 0
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    total_chars += len(text)
                    cyrillic_chars += sum(1 for c in text if '\u0400' <= c <= '\u04FF')
        
        if total_chars > 0:
            cyrillic_ratio = cyrillic_chars / total_chars
            
            if cyrillic_ratio < 0.3:
                self.errors.append(f"Недостаточно кириллицы: {cyrillic_ratio:.1%} (минимум 30%)")
            else:
                self.info.append(f"✅ Кириллица: {cyrillic_ratio:.1%}")
        
        return len(self.errors) == 0
    
    def verify_notes(self) -> bool:
        """Проверка заметок докладчика"""
        print("🔍 Проверка заметок докладчика...")
        
        slides_with_notes = 0
        total_slides = len(self.prs.slides)
        
        for slide in self.prs.slides:
            if slide.notes_slide.notes_text_frame.text.strip():
                slides_with_notes += 1
        
        notes_ratio = slides_with_notes / total_slides if total_slides > 0 else 0
        
        if notes_ratio < 0.5:
            self.warnings.append(f"Мало заметок докладчика: {notes_ratio:.1%} слайдов (рекомендуется 80%)")
        elif notes_ratio < 0.8:
            self.warnings.append(f"Недостаточно заметок докладчика: {notes_ratio:.1%} слайдов (рекомендуется 80%)")
        else:
            self.info.append(f"✅ Заметки докладчика: {notes_ratio:.1%} слайдов")
        
        return notes_ratio >= 0.5
    
    def verify_figures(self) -> bool:
        """Проверка рисунков"""
        print("🔍 Проверка рисунков...")
        
        figure_count = 0
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    figure_count += 1
        
        if figure_count < 1:
            self.warnings.append("Нет рисунков в презентации")
        else:
            self.info.append(f"✅ Рисунков: {figure_count}")
        
        return True
    
    def verify_text_length(self) -> bool:
        """Проверка длины текста на слайдах"""
        print("🔍 Проверка длины текста...")
        
        long_slides = []
        
        for i, slide in enumerate(self.prs.slides, 1):
            line_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    line_count += len(text.split('\n'))
            
            if line_count > 10:
                long_slides.append((i, line_count))
        
        if long_slides:
            for slide_num, lines in long_slides:
                self.warnings.append(f"Слайд {slide_num}: слишком много текста ({lines} строк, рекомендуется до 7)")
        else:
            self.info.append("✅ Длина текста на слайдах в норме")
        
        return True
    
    def verify_font_size(self) -> bool:
        """Проверка размера шрифта"""
        print("🔍 Проверка размера шрифта...")
        
        small_fonts = []
        
        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size.pt < 18:
                                small_fonts.append((i, run.font.size.pt))
        
        if small_fonts:
            for slide_num, size in small_fonts[:3]:  # Показываем первые 3
                self.warnings.append(f"Слайд {slide_num}: мелкий шрифт ({size:.0f}pt, рекомендуется минимум 18pt)")
        else:
            self.info.append("✅ Размер шрифта в норме")
        
        return True
    
    def verify_aspect_ratio(self) -> bool:
        """Проверка соотношения сторон"""
        print("🔍 Проверка соотношения сторон...")
        
        width = self.prs.slide_width
        height = self.prs.slide_height
        ratio = width / height
        
        # 16:9 = 1.778
        if abs(ratio - 1.778) < 0.01:
            self.info.append(f"✅ Формат: 16:9")
        # 4:3 = 1.333
        elif abs(ratio - 1.333) < 0.01:
            self.warnings.append(f"Формат 4:3 (рекомендуется 16:9)")
        else:
            self.warnings.append(f"Нестандартный формат: {ratio:.2f} (рекомендуется 16:9)")
        
        return True
    
    def verify_gost_compliance(self) -> bool:
        """Проверка соответствия ГОСТ (базовая)"""
        print("🔍 Проверка соответствия ГОСТ...")
        
        # Проверка десятичных разделителей (должна быть запятая, а не точка)
        decimal_point_issues = []
        
        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    # Ищем числа с точкой (например, 3.14)
                    matches = re.findall(r'\d+\.\d+', text)
                    if matches:
                        decimal_point_issues.append((i, matches))
        
        if decimal_point_issues:
            for slide_num, matches in decimal_point_issues[:3]:
                self.warnings.append(f"Слайд {slide_num}: десятичная точка вместо запятой ({matches[0]})")
        else:
            self.info.append("✅ Десятичные разделители в норме")
        
        return True
    
    def run_all_checks(self) -> dict:
        """Запустить все проверки"""
        print("\n" + "="*60)
        print("ВЕРИФИКАЦИЯ ПРЕЗЕНТАЦИИ")
        print("="*60 + "\n")
        
        checks = [
            self.verify_structure,
            self.verify_language,
            self.verify_notes,
            self.verify_figures,
            self.verify_text_length,
            self.verify_font_size,
            self.verify_aspect_ratio,
            self.verify_gost_compliance
        ]
        
        for check in checks:
            check()
            print()
        
        # Итоговый отчет
        print("="*60)
        print("ИТОГОВЫЙ ОТЧЕТ")
        print("="*60 + "\n")
        
        if self.errors:
            print("❌ ОШИБКИ:")
            for error in self.errors:
                print(f"  • {error}")
            print()
        
        if self.warnings:
            print("⚠️  ПРЕДУПРЕЖДЕНИЯ:")
            for warning in self.warnings:
                print(f"  • {warning}")
            print()
        
        if self.info:
            print("ℹ️  ИНФОРМАЦИЯ:")
            for info in self.info:
                print(f"  • {info}")
            print()
        
        # Общая оценка
        total_issues = len(self.errors) + len(self.warnings)
        
        if len(self.errors) == 0:
            if len(self.warnings) == 0:
                print("✅ ОТЛИЧНО: Презентация соответствует всем требованиям")
                quality = "excellent"
            elif len(self.warnings) <= 3:
                print("✅ ХОРОШО: Презентация соответствует основным требованиям")
                quality = "good"
            else:
                print("⚠️  УДОВЛЕТВОРИТЕЛЬНО: Есть замечания, требующие внимания")
                quality = "satisfactory"
        else:
            print("❌ НЕУДОВЛЕТВОРИТЕЛЬНО: Есть критические ошибки")
            quality = "poor"
        
        print("="*60 + "\n")
        
        return {
            'quality': quality,
            'errors': self.errors,
            'warnings': self.warnings,
            'info': self.info,
            'total_issues': total_issues
        }
    
    def save_report(self, output_path: str):
        """Сохранить отчет в файл"""
        report = self.run_all_checks()
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("ОТЧЕТ О ВЕРИФИКАЦИИ ПРЕЗЕНТАЦИИ\n")
            f.write("="*60 + "\n\n")
            f.write(f"Файл: {self.pptx_path}\n")
            f.write(f"Качество: {report['quality']}\n")
            f.write(f"Всего замечаний: {report['total_issues']}\n\n")
            
            if report['errors']:
                f.write("ОШИБКИ:\n")
                for error in report['errors']:
                    f.write(f"  • {error}\n")
                f.write("\n")
            
            if report['warnings']:
                f.write("ПРЕДУПРЕЖДЕНИЯ:\n")
                for warning in report['warnings']:
                    f.write(f"  • {warning}\n")
                f.write("\n")
            
            if report['info']:
                f.write("ИНФОРМАЦИЯ:\n")
                for info in report['info']:
                    f.write(f"  • {info}\n")
                f.write("\n")
        
        print(f"📄 Отчет сохранен: {output_path}")


def main():
    """Пример использования"""
    import sys
    
    if len(sys.argv) < 2:
        print("Использование: python verify_presentation.py <путь_к_pptx>")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    
    if not Path(pptx_path).exists():
        print(f"❌ Файл не найден: {pptx_path}")
        sys.exit(1)
    
    verifier = PresentationVerifier(pptx_path)
    report = verifier.run_all_checks()
    
    # Сохранить отчет
    output_path = Path(pptx_path).parent / "verification_report.txt"
    verifier.save_report(str(output_path))


if __name__ == "__main__":
    main()

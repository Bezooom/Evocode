#!/usr/bin/env python3
"""
Пример создания презентации для защиты диссертации
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class DissertationPresentation:
    """Создание презентации для защиты диссертации"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)  # 16:9
        self.prs.slide_height = Inches(5.625)
    
    def add_title_slide(self, title: str, author: str, specialty: str, 
                       university: str, year: int):
        """Титульный слайд для защиты диссертации"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Университет (вверху)
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.5)
        
        uni_box = slide.shapes.add_textbox(left, top, width, height)
        uni_frame = uni_box.text_frame
        uni_frame.text = university
        uni_frame.paragraphs[0].font.size = Pt(16)
        uni_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Название диссертации (центр)
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True
        
        # Специальность
        left = Inches(1)
        top = Inches(3.2)
        width = Inches(8)
        height = Inches(0.4)
        
        spec_box = slide.shapes.add_textbox(left, top, width, height)
        spec_frame = spec_box.text_frame
        spec_frame.text = f"Специальность: {specialty}"
        spec_frame.paragraphs[0].font.size = Pt(18)
        spec_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Автор
        left = Inches(1)
        top = Inches(4)
        width = Inches(8)
        height = Inches(0.8)
        
        author_box = slide.shapes.add_textbox(left, top, width, height)
        author_frame = author_box.text_frame
        author_frame.text = f"Диссертант: {author}\n{year} г."
        author_frame.paragraphs[0].font.size = Pt(20)
        author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_relevance_slide(self, relevance_items: list, notes: str = ""):
        """Слайд: Актуальность темы диссертации"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Актуальность темы диссертации"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, item in enumerate(relevance_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(22)
            p.level = 0
            p.space_before = Pt(6)
        
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def add_goal_and_tasks_slide(self, goal: str, tasks: list, notes: str = ""):
        """Слайд: Цель и задачи диссертационного исследования"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Цель и задачи диссертационного исследования"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Цель
        p = text_frame.paragraphs[0]
        p.text = f"Цель: {goal}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_after = Pt(12)
        
        # Задачи
        p = text_frame.add_paragraph()
        p.text = "Задачи:"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_before = Pt(6)
        
        for i, task in enumerate(tasks, 1):
            p = text_frame.add_paragraph()
            p.text = f"{i}. {task}"
            p.font.size = Pt(20)
            p.level = 1
            p.space_before = Pt(4)
        
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def add_novelty_slide(self, novelty_items: list, notes: str = ""):
        """Слайд: Научная новизна"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Научная новизна"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, item in enumerate(novelty_items, 1):
            if i == 1:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"{i}. {item}"
            p.font.size = Pt(20)
            p.level = 0
            p.space_before = Pt(8)
        
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def add_significance_slide(self, theoretical: list, practical: list, notes: str = ""):
        """Слайд: Теоретическая и практическая значимость"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Теоретическая и практическая значимость"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Теоретическая значимость
        p = text_frame.paragraphs[0]
        p.text = "Теоретическая значимость:"
        p.font.size = Pt(22)
        p.font.bold = True
        
        for item in theoretical:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(4)
        
        # Практическая значимость
        p = text_frame.add_paragraph()
        p.text = "Практическая значимость:"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_before = Pt(12)
        
        for item in practical:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(4)
        
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def add_approbation_slide(self, conferences: list, publications: list, notes: str = ""):
        """Слайд: Апробация результатов"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Апробация результатов исследования"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Конференции
        p = text_frame.paragraphs[0]
        p.text = "Доклады на конференциях:"
        p.font.size = Pt(22)
        p.font.bold = True
        
        for conf in conferences:
            p = text_frame.add_paragraph()
            p.text = f"• {conf}"
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(4)
        
        # Публикации
        p = text_frame.add_paragraph()
        p.text = f"Публикации: {len(publications)} статей"
        p.font.size = Pt(22)
        p.font.bold = True
        p.space_before = Pt(12)
        
        for pub in publications[:3]:  # Показываем первые 3
            p = text_frame.add_paragraph()
            p.text = f"• {pub}"
            p.font.size = Pt(16)
            p.level = 1
            p.space_before = Pt(4)
        
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def add_publications_slide(self, publications: list, notes: str = ""):
        """Слайд: Публикации по теме диссертации"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Публикации по теме диссертации"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, pub in enumerate(publications, 1):
            if i == 1:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"{i}. {pub}"
            p.font.size = Pt(16)
            p.level = 0
            p.space_before = Pt(6)
        
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    def save(self, output_path: str):
        """Сохранить презентацию"""
        self.prs.save(output_path)
        print(f"✅ Презентация для защиты диссертации создана: {output_path}")


def main():
    """Пример создания презентации для защиты кандидатской диссертации"""
    
    pres = DissertationPresentation()
    
    # 1. Титульный слайд
    pres.add_title_slide(
        title="Молекулярные механизмы регуляции клеточного цикла белком X",
        author="Иванов Иван Иванович",
        specialty="03.01.03 — Молекулярная биология",
        university="Московский государственный университет имени М.В. Ломоносова",
        year=2024
    )
    
    # 2. Актуальность
    pres.add_relevance_slide(
        relevance_items=[
            "Нарушения клеточного цикла лежат в основе онкологических заболеваний",
            "Белок X является ключевым регулятором, но его механизм действия не изучен",
            "Понимание молекулярных механизмов необходимо для разработки таргетной терапии",
            "Отсутствуют данные о взаимодействии белка X с другими регуляторами"
        ],
        notes="Актуальность темы диссертации обусловлена тем, что нарушения клеточного цикла лежат в основе большинства онкологических заболеваний. Белок X является ключевым регулятором клеточного цикла, однако молекулярные механизмы его действия остаются недостаточно изученными. Понимание этих механизмов необходимо для разработки новых подходов к таргетной терапии рака. Кроме того, в литературе отсутствуют данные о взаимодействии белка X с другими регуляторами клеточного цикла."
    )
    
    # 3. Цель и задачи
    pres.add_goal_and_tasks_slide(
        goal="Изучить молекулярные механизмы регуляции клеточного цикла белком X",
        tasks=[
            "Идентифицировать белки-партнеры белка X методом масс-спектрометрии",
            "Исследовать влияние белка X на прогрессию клеточного цикла",
            "Определить сайты фосфорилирования белка X и их функциональную роль",
            "Изучить клиническое значение экспрессии белка X при раке"
        ],
        notes="Основная цель диссертационного исследования — изучить молекулярные механизмы регуляции клеточного цикла белком X. Для достижения этой цели были поставлены четыре задачи. Первая задача — идентифицировать белки-партнеры белка X с использованием современных протеомных подходов. Вторая задача — исследовать влияние белка X на прогрессию клеточного цикла в различных клеточных линиях. Третья задача — определить сайты посттрансляционных модификаций белка X и их функциональную роль. Четвертая задача — изучить клиническое значение экспрессии белка X при различных типах рака."
    )
    
    # 4. Научная новизна
    pres.add_novelty_slide(
        novelty_items=[
            "Впервые идентифицированы 15 новых белков-партнеров белка X",
            "Установлена роль белка X в регуляции перехода G1/S фазы клеточного цикла",
            "Обнаружены 3 новых сайта фосфорилирования, критичных для функции белка X",
            "Показана корреляция между экспрессией белка X и прогнозом при раке молочной железы"
        ],
        notes="Научная новизна работы заключается в следующем. Во-первых, впервые с использованием масс-спектрометрии высокого разрешения идентифицированы 15 новых белков-партнеров белка X. Во-вторых, установлена ключевая роль белка X в регуляции перехода из G1 в S фазу клеточного цикла. В-третьих, обнаружены три новых сайта фосфорилирования, которые являются критичными для функции белка X. В-четвертых, впервые показана корреляция между уровнем экспрессии белка X и клиническим прогнозом при раке молочной железы."
    )
    
    # 5. Теоретическая и практическая значимость
    pres.add_significance_slide(
        theoretical=[
            "Расширены представления о молекулярных механизмах регуляции клеточного цикла",
            "Установлена новая сигнальная ось белок X — белок Y — CDK2",
            "Предложена модель регуляции перехода G1/S фазы"
        ],
        practical=[
            "Белок X может служить биомаркером прогноза при раке молочной железы",
            "Идентифицированы потенциальные мишени для таргетной терапии",
            "Разработан метод количественной оценки активности белка X"
        ],
        notes="Теоретическая значимость работы состоит в том, что получены новые фундаментальные знания о молекулярных механизмах регуляции клеточного цикла. Установлена новая сигнальная ось, включающая белок X, белок Y и циклин-зависимую киназу CDK2. Предложена комплексная модель регуляции перехода из G1 в S фазу клеточного цикла. Практическая значимость заключается в том, что белок X может использоваться в качестве биомаркера прогноза при раке молочной железы. Идентифицированы потенциальные мишени для разработки новых противоопухолевых препаратов. Разработан и валидирован метод количественной оценки активности белка X, который может применяться в клинической практике."
    )
    
    # 6. Апробация
    pres.add_approbation_slide(
        conferences=[
            "Международная конференция по молекулярной биологии (Москва, 2023)",
            "European Cancer Congress (Париж, 2023)",
            "Съезд биохимиков России (Санкт-Петербург, 2024)"
        ],
        publications=[
            "Иванов И.И. и др. Роль белка X в регуляции клеточного цикла // Молекулярная биология. 2023. Т. 57. № 3. С. 450-465.",
            "Ivanov I.I. et al. Protein X regulates G1/S transition // Nature Cell Biology. 2024. Vol. 26. P. 123-135.",
            "Иванов И.И., Петров П.П. Фосфорилирование белка X // Биохимия. 2024. Т. 89. № 2. С. 234-245."
        ],
        notes="Результаты диссертационного исследования были представлены на трех крупных научных конференциях, включая международную конференцию по молекулярной биологии в Москве и Европейский онкологический конгресс в Париже. По теме диссертации опубликовано 3 статьи в рецензируемых научных журналах, в том числе 2 статьи в журналах, входящих в перечень ВАК, и 1 статья в высокорейтинговом международном журнале Nature Cell Biology."
    )
    
    # 7. Публикации (детально)
    pres.add_publications_slide(
        publications=[
            "Иванов И.И., Петров П.П., Сидоров С.С. Роль белка X в регуляции клеточного цикла // Молекулярная биология. 2023. Т. 57. № 3. С. 450-465. (ВАК, Scopus)",
            "Ivanov I.I., Petrov P.P., Sidorov S.S. Protein X regulates G1/S transition through CDK2 phosphorylation // Nature Cell Biology. 2024. Vol. 26. P. 123-135. (IF=20.042, Q1)",
            "Иванов И.И., Петров П.П. Фосфорилирование белка X регулирует его взаимодействие с белком Y // Биохимия. 2024. Т. 89. № 2. С. 234-245. (ВАК, Scopus)"
        ],
        notes="Все публикации по теме диссертации размещены в рецензируемых научных журналах. Первая статья опубликована в журнале Молекулярная биология, который входит в перечень ВАК и индексируется в базе данных Scopus. Вторая статья опубликована в высокорейтинговом международном журнале Nature Cell Biology с импакт-фактором более 20. Третья статья опубликована в журнале Биохимия, также входящем в перечень ВАК и Scopus."
    )
    
    # Сохранить
    pres.save("dissertation_presentation_ru.pptx")
    print(f"🎓 Презентация для защиты диссертации готова!")


if __name__ == "__main__":
    main()
